import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup
import requests
from io import BytesIO

def html_to_pptx(html_file, output_file):
    prs = Presentation()
    
    with open(html_file, 'r', encoding='utf-8') as file:
        html_content = file.read()
    
    soup = BeautifulSoup(html_content, 'html.parser')
    
    title = soup.title.string if soup.title else "Converted Presentation"
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    main_content = soup.find('div', class_='main-container')
    
    if main_content:
        for element in main_content.find_all(['h1', 'h2', 'h3', 'p', 'ul', 'ol', 'img']):
            if element.name == 'img':
                img_slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(img_slide_layout)
                
                img_src = element.get('src')
                if img_src.startswith('http'):
                    response = requests.get(img_src)
                    img_stream = BytesIO(response.content)
                else:
                    img_stream = img_src  # Local file path
                
                slide.shapes.add_picture(img_stream, Inches(1), Inches(1), width=Inches(8), height=Inches(5.5))
            else:
                content_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(content_slide_layout)
                
                if element.name in ['h1', 'h2', 'h3']:
                    title_shape = slide.shapes.title
                    title_shape.text = element.get_text(strip=True)
                    title_shape.text_frame.paragraphs[0].font.size = Pt(40 if element.name == 'h1' else 36 if element.name == 'h2' else 32)
                    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
                content_shape = slide.shapes.placeholders[1]
                text_frame = content_shape.text_frame
                text_frame.clear()  # Clear default text
                
                if element.name in ['p', 'ul', 'ol']:
                    if element.name == 'p':
                        p = text_frame.add_paragraph()
                        p.text = element.get_text(strip=True)
                        p.font.size = Pt(18)
                    else:  # ul or ol
                        for li in element.find_all('li'):
                            item = text_frame.add_paragraph()
                            item.text = '• ' + li.get_text(strip=True)
                            item.level = 1
                            item.font.size = Pt(18)
                
                for paragraph in text_frame.paragraphs:
                    paragraph.font.color.rgb = RGBColor(0, 0, 0)
    
    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

# Example usage
html_file = 'thesis-bachelor-2024.html'
output_file = 'output_from_python.pptx'
html_to_pptx(html_file, output_file)